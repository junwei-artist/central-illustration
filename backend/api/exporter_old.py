from fastapi import APIRouter, HTTPException, Depends, Query
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from db.session import get_db
from models.demonstration import Demonstration
from core.process_manager import process_manager
import httpx
from pathlib import Path
import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import io
import base64
from typing import Literal
import os

router = APIRouter()

async def capture_slide_screenshots(demo_url: str, num_slides: int) -> list[bytes]:
    """Capture screenshots of each slide from the hero section"""
    screenshots = []
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page()
        
        # Set viewport for 16:9 or 4:3 based on format
        await page.set_viewport_size({"width": 1920, "height": 1080})
        
        # Navigate to the demo URL
        await page.goto(demo_url, wait_until='networkidle', timeout=30000)
        
        # Wait for content to load
        await asyncio.sleep(2)
        
        for i in range(num_slides):
            try:
                # Click on the numbered button to navigate to the slide
                slide_button = page.locator(f'button:has-text("{i + 1}")')
                if await slide_button.count() > 0:
                    await slide_button.click()
                    await asyncio.sleep(1)  # Wait for transition
                
                # Capture screenshot of the hero section
                screenshot = await page.screenshot(full_page=False)
                screenshots.append(screenshot)
            except Exception as e:
                print(f"Error capturing slide {i}: {e}")
                # Take screenshot anyway
                screenshot = await page.screenshot(full_page=False)
                screenshots.append(screenshot)
        
        await browser.close()
    
    return screenshots


def create_powerpoint_presentation(screenshots: list[bytes], aspect_ratio: str) -> bytes:
    """Create a PowerPoint presentation from screenshots"""
    prs = Presentation()
    
    # Set slide dimensions
    if aspect_ratio == '16:9':
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:  # 4:3
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    for screenshot in screenshots:
        # Create a slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Convert IPO screenshot to Image
        img = Image.open(io.BytesIO(screenshot))
        
        # Resize image to fit slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        img.thumbnail((int(slide_width), int(slide_height)), Image.Resampling.LANCZOS)
        
        # Save to buffer
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        
        # Add image to slide
        left = (slide_width - Inches(img.width / 96)) / 2
        top = (slide_height - Inches(img.height / 96)) / 2
        slide.shapes.add_picture(img_buffer, left, top, 
                                width=Inches(img.width / 96), 
                                height=Inches(img.height / 96))
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return buffer.getvalue()


async def create_pdf_from_screenshots(screenshots: list[bytes], aspect_ratio: str) -> bytes:
    """Create a PDF from screenshots using Playwright"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    buffer = io.BytesIO()
    
    if aspect_ratio == '16:9':
        page_width, page_height = 1920, 1080
    else:  # 4:3
        page_width, page_height = 1920, 1440
    
    # Use Playwright to create PDF
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page()
        
        # Create HTML with all screenshots
        html = f"<div style='width: {page_width}px; height: {page_height}px;'>"
        for screenshot in screenshots:
            b64 = base64.b64encode(screenshot).decode()
            html += f"<div style='page-break-after: always; width: {page_width}px; height: {page_height}px;'>"
            html += f"<img src='data:image/png;base64,{b64}' style='width: 100%; height: 100%; object-fit: contain;' />"
            html += "</div>"
        html += "</div>"
        
        await page.set_content(html)
        pdf_bytes = await page.pdf(format='A4' if aspect_ratio == '4:3' else None, width=f'{page_width}px', height=f'{page_height}px')
        
        await browser.close()
        
    return pdf_bytes


@router.post("/export/{demo_id}")
async def export_demo(
    demo_id: int,
    format: Literal['ppt_169', 'ppt_43', 'pdf_169', 'pdf_43'] = Query(...),
    db: Session = Depends(get_db)
):
    """
    Export a demo's hero section to PowerPoint or PDF
    """
    # Get the demo
    demo = db.query(Demonstration).filter(Demonstration.id == demo_id).first()
    if not demo:
        raise HTTPException(status_code=404, detail="Demo not found")
    
    # Check if demo is running
    status_info = process_manager.get_demo_status(demo.folder_name)
    
    if status_info['status'] != 'running' or not status_info.get('port'):
        raise HTTPException(status_code=503, detail="Demo service is not running")
    
    # Get the URL
    port = status_info['port']
    demo_url = f"http://127.0.0.1:{port}"
    
    try:
        # Parse format
        is_pdf = format.startswith('pdf')
        aspect_ratio = format.split('_')[1]
        
        # Number of slides (assuming 8 slides for now, can be dynamic)
        num_slides = 8
        
        # Capture screenshots
        screenshots = await capture_slide_screenshots(demo_url, num_slides)
        
        if not screenshots:
            raise HTTPException(status_code=500, detail="Failed to capture slides")
        
        # Generate export file
        if is_pdf:
            output = await create_pdf_from_screenshots(screenshots, aspect_ratio)
            filename = f"{demo.folder_name}_export.pdf"
            media_type = "application/pdf"
        else:
            output = create_powerpoint_presentation(screenshots, aspect_ratio)
            filename = f"{demo.folder_name}_export.pptx"
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        # Return file
        return StreamingResponse(
            io.BytesIO(output),
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={filename}"
            }
        )
        
    except Exception as e:
        print(f"Export error: {e}")
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

