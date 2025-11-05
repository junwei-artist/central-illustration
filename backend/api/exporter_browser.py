from fastapi import APIRouter, HTTPException, Depends, Query
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from db.session import get_db
from models.demonstration import Demonstration
from core.process_manager import process_manager
import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from typing import Literal, List, Dict, Any
import re

router = APIRouter()


async def extract_slide_data(demo_url: str, num_slides: int) -> List[Dict[str, Any]]:
    """Extract actual content from each slide including text, styles, and structure"""
    slides_data = []
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page()
        await page.set_viewport_size({"width": 1920, "height": 1080})
        await page.goto(demo_url, wait_until='networkidle', timeout=30000)
        await asyncio.sleep(2)
        
        for i in range(num_slides):
            try:
                # Navigate to slide
                slide_button = page.locator(f'button:has-text("{i + 1}")')
                if await slide_button.count() > 0:
                    await slide_button.click()
                    await asyncio.sleep(2)  # Wait longer for transitions
                
                # Extract slide content - need to find the ACTIVE slide specifically
                slide_data = await page.evaluate(f"""
                    (slideIndex) => {{
                        const section = document.querySelector('section');
                        if (!section) return {{}};
                        
                        // Find the active motion.div for this slide
                        const allSlides = section.querySelectorAll('div[class*="absolute"]');
                        let activeSlide = null;
                        
                        for (const div of allSlides) {{
                            const opacity = window.getComputedStyle(div).opacity;
                            const visibility = window.getComputedStyle(div).visibility;
                            // Look for visible slide (not pointer-events-none invisible)
                            if (parseFloat(opacity) > 0.5 && visibility !== 'hidden' && !div.classList.contains('pointer-events-none')) {{
                                activeSlide = div;
                                break;
                            }}
                        }}
                        
                        if (!activeSlide) return {{}};
                        
                        const getTextFrom = (selector, parent) => {{
                            const el = parent ? parent.querySelector(selector) : section.querySelector(selector);
                            return el ? el.textContent.trim() : '';
                        }};
                        
                        const getAllTextFrom = (selector, parent) => {{
                            const elements = parent ? parent.querySelectorAll(selector) : section.querySelectorAll(selector);
                            return Array.from(elements).map(el => el.textContent.trim()).filter(Boolean);
                        }};
                        
                        return {{
                            title: getTextFrom('h1', activeSlide),
                            subtitle: getAllTextFrom('h1', activeSlide)[1] || '',
                            description: getTextFrom('p', activeSlide),
                            points: getAllTextFrom('li', activeSlide),
                            bgColor: window.getComputedStyle(activeSlide).backgroundImage || window.getComputedStyle(activeSlide).backgroundColor,
                            textColor: window.getComputedStyle(activeSlide.querySelector('h1') || activeSlide).color,
                            hasSvg: activeSlide.querySelector('svg') !== null
                        }};
                    }}
                """, i)
                
                print(f"Slide {i} data: {slide_data}")
                slides_data.append(slide_data)
            except Exception as e:
                print(f"Error extracting slide {i}: {e}")
                import traceback
                traceback.print_exc()
                slides_data.append({})
        
        await browser.close()
    
    return slides_data


def parse_gradient_color(gradient_str: str) -> str:
    """Parse CSS gradient to get a solid color"""
    # Try to extract a color from gradient
    match = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', gradient_str)
    if match:
        return f"#{int(match.group(1)):02x}{int(match.group(2)):02x}{int(match.group(3)):02x}"
    
    # Try hex color
    match = re.search(r'#([0-9a-fA-F]{6})', gradient_str)
    if match:
        return f"#{match.group(1)}"
    
    return "#3B82F6"  # Default blue


def create_powerpoint_with_objects(slides_data: List[Dict], aspect_ratio: str) -> bytes:
    """Create PowerPoint with actual text boxes and shapes instead of screenshots"""
    prs = Presentation()
    
    if aspect_ratio == '16:9':
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    for data in slides_data:
        if not data:
            continue
            
        # Create blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Set background
        if data.get('bgColor'):
            try:
                background = slide.background
                fill = background.fill
                fill.solid()
                color = parse_gradient_color(data['bgColor'])
                # Convert hex to RGB
                if color.startswith('#'):
                    r = int(color[1:3], 16)
                    g = int(color[3:5], 16)
                    b = int(color[5:7], 16)
                    fill.fore_color.rgb = RGBColor(r, g, b)
                else:
                    fill.fore_color.rgb = RGBColor(59, 130, 246)  # Default blue
            except Exception as e:
                print(f"Background color error: {e}")
                pass
        
        # Add title
        if data.get('title'):
            left = Inches(1)
            top = Inches(1.5)
            width = prs.slide_width - Inches(2)
            height = Inches(1.5)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = data['title']
            
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(60)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add subtitle
        if data.get('subtitle'):
            left = Inches(1)
            top = Inches(3)
            width = prs.slide_width - Inches(2)
            height = Inches(0.8)
            
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = data['subtitle']
            
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.bold = False
            subtitle_para.font.color.rgb = RGBColor(220, 220, 220)
        
        # Add description
        if data.get('description'):
            left = Inches(1)
            top = Inches(4)
            width = prs.slide_width - Inches(2)
            height = Inches(1)
            
            desc_box = slide.shapes.add_textbox(left, top, width, height)
            desc_frame = desc_box.text_frame
            desc_frame.text = data['description']
            
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(18)
            desc_para.font.color.rgb = RGBColor(200, 200, 200)
        
        # Add bullet points
        if data.get('points'):
            left = Inches(1)
            top = Inches(5.5)
            width = prs.slide_width - Inches(2)
            max_height = prs.slide_height - Inches(5.5)
            
            points_box = slide.shapes.add_textbox(left, top, width, max_height)
            points_frame = points_box.text_frame
            points_frame.word_wrap = True
            
            for i, point in enumerate(data['points']):
                if i > 0:
                    points_frame.add_paragraph()
                para = points_frame.paragraphs[i]
                para.text = point.replace('—', '-').replace('•', '')
                para.level = 0
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(240, 240, 240)
                para.space_after = Pt(4)
        
        # Add decorative shape to mimic SVG presence
        if data.get('hasSvg'):
            # Add a simple circle or rectangle as decoration
            left = Inches(prs.slide_width.inches - 3)
            top = Inches(2)
            width = Inches(2)
            height = Inches(2)
            
            shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.fill.transparency = 0.8
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@router.post("/export/{demo_id}")
async def export_demo(
    demo_id: int,
    format: Literal['ppt_169', 'ppt_43', 'pdf_169', 'pdf_43'] = Query(...),
    db: Session = Depends(get_db)
):
    """Export demo to PowerPoint or PDF with actual content objects"""
    demo = db.query(Demonstration).filter(Demonstration.id == demo_id).first()
    if not demo:
        raise HTTPException(status_code=404, detail="Demo not found")
    
    status_info = process_manager.get_demo_status(demo.folder_name)
    if status_info['status'] != 'running' or not status_info.get('port'):
        raise HTTPException(status_code=503, detail="Demo service is not running")
    
    port = status_info['port']
    demo_url = f"http://127.0.0.1:{port}"
    
    try:
        is_pdf = format.startswith('pdf')
        aspect_ratio = format.split('_')[1]
        num_slides = 8
        
        # Extract actual slide data
        slides_data = await extract_slide_data(demo_url, num_slides)
        
        print(f"Extracted {len(slides_data)} slides")
        
        if not slides_data:
            raise HTTPException(status_code=500, detail="Failed to extract slide data")
        
        # For now, PDF will still use screenshots approach (can be enhanced later)
        if is_pdf:
            raise HTTPException(status_code=501, detail="PDF export with real objects not yet implemented")
        
        # Create PowerPoint with real objects
        output = create_powerpoint_with_objects(slides_data, aspect_ratio)
        filename = f"{demo.folder_name}_export.pptx"
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        return StreamingResponse(
            io.BytesIO(output),
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Export error: {e}")
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

