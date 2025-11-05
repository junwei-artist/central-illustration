from fastapi import APIRouter, HTTPException, Depends, Query
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from db.session import get_db
from models.demonstration import Demonstration
from pptx import Presentation
from pptx.util import Inches, greediness Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import io
import json
import re
import ast
from typing import Literal, List, Dict, Any

router = APIRouter()


def extract_hero_slides_from_project(project_path: Path) -> List[Dict]:
    """Extract heroSlides array from a Next.js project's page.tsx"""
    page_file = project_path / 'app' / 'page.tsx'
    
    if not page_file.exists():
        return []
    
    with open(page_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Find the heroSlides array using regex
    match = re.search(r'const heroSlides = (\[[\s\S]*?\])', content)
    if not match:
        return []
    
    # Try to parse as JavaScript array
    js_array = match.group(1)
    
    # Convert JS syntax to Python dict
    # Replace single quotes with double quotes
    js_array = js_array.replace("'", '"')
    # Handle JS array syntax
    js_array = js_array.replace('[', '(').replace(']', ')')
    
    # Parse line by line for safety
    slides = []
    current_slide = {}
    
    # Extract slide objects manually
    slide_matches = re.finditer(r'\{[^}]*title[^}]*\}', match.group(1), re.DOTALL)
    for m in slide_matches:
        slide_str = m.group(0)
        
        # Extract individual fields
        title_match = re.search(r"'title':\s*'(.*?)'", slide_str)
        subtitle_match = re.search(r"'subtitle':\s*'(.*?)'", slide_str)
        desc_match = re.search(r"'description':\s*'(.*?)'", slide_str)
        colors_match = re.search(r"'colors':\s*'(.*?)'", slide_str)
        
        # Extract points array
        points = []
        points_match = re.search(r"'points':\s*\[(.*?)\]", slide_str, re.DOTALL)
        if points_match:
            points_text = points_match.group(1)
            point_matches = re.finditer(r"'([^']+)'", points_text)
            for p in point_matches:
                points.append(p.group(1))
        
        slide = {
            'title': title_match.group(1) if title_match else '',
            'subtitle': subtitle_match.group(1) if subtitle_match else '',
            'description': desc_match.group(1) if desc_match else '',
            'colors': colors_match.group(1) if colors_match else '',
            'points': points
        }
        slides.append(slide)
    
    return slides


def parse_gradient_colors(gradient_str: str) -> tuple:
    """Parse Tailwind gradient classes to get RGB colors"""
    default_colors = {
        'from-blue-600 to-purple-700': (RGBColor(37, 99, 235), RGBColor(126, 34, 206)),
        'from-indigo-600 to-cyan-600': (RGBColor(79, 70, 229), RGBColor(8, 145, 178)),
        'from-fuchsia-600 to-roseждав 600': (RGBColor(192, 38, 211), RGBColor(225, 29, 72)),
        'from-emerald-600 to-teal-600': (RGBColor(5, 150, 105), RGBColor(13, 148, 136)),
        'from-amber-600 to-orange-600': (RGBColor(217, 119, 6), RGBColor(234, 88, 12)),
        'from-violet-600 to-blue-600': (RGBColor(124, 58, 237), RGBColor(37, 99, 235)),
        'from-sky-600 to-cyan-600': (RGBColor(2, 132, 199), RGBColor(8, 145, 178)),
        'from-slate-700 to-gray-800': (RGBColor(51, 65, 85), RGBColor(31, 41, 55)),
    }
    return default_colors.get(gradient_str, (RGBColor(59, 130, 246), RGBColor(147, 51, 234)))


def create_powerpoint_from_slides(slides_data: List[Dict], aspect_ratio: str) -> bytes:
    """Create PowerPoint from generic slide data"""
    prs = Presentation()
    
    if aspect_ratio == '16:9':
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    for i, data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        colors = parse_gradient_colors(data.get('colors', ''))
        background = slide.background
        fill = background.fill
        fill.solid()
rect.fore_color.rgb = colors[0]
        
        # Add title
        if data.get('title'):
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), prs.slide_width - Inches(2), Inches(1.5))
            title_box.text_frame.text = data['title']
            title_box.text_frame.paragraphs[0].font.size = Pt(56)
            title_box.text_frame.paragraphs[0].font.bold = True
            title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add subtitle
        if data.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), prs.slide_width - Inches(2), Inches(0.8))
            subtitle_box.text_frame.text = data['subtitle']
            subtitle_box.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(240, 240, 240)
        
        # Add description
        if data.get('description'):
            desc_box = slide.shapes.add_textbox(Inches(1), Inches(3.6), prs.slide_width - Inches(2), Inches(1))
            desc_box.text_frame.text = data['description']
            desc_box.text_frame.paragraphs[0].font.size = Pt(16)
            desc_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
        
        # Add bullet points
        if data.get('points'):
            points_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), prs.slide_width - Inches(2.5), prs.slide_height - Inches(4.8))
            points_frame = points_box.text_frame
            points_frame.word_wrap = True
            
            for j, point in enumerate(data['points']):
                if j > 0:
                    points_frame.add_paragraph()
                para = points_frame.paragraphs[j]
                para.text = point.replace('—', '-').replace('•', '').strip()
                para.level = 0
                para.font.size = Pt(13)
                para.font.color.rgb = RGBColor(240, 240, 240)
                para.space_after = Pt(3)
        
        # Add simple decorative shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, prs.slide_width - Inches(3.5), Inches(1.5), Inches(2.5), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.fill.transparency = 0.15
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(2)
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@router.post("/export/{demo_id}")
async def export_demo(
    demo_id: int,
    format: Literal['ppt_169', 'ppt_43', 'pdf_169', 'pdf_43'] = Query(...),
    db: Session = Depends(get_db)
):
    """Generic export that works for any project's hero section"""
    demo = db.query(Demonstration).filter(Demonstration.id == demo_id).first()
    if not demo:
        raise HTTPException(status_code=404, detail="Demo not found")
    
    try:
        is_pdf = format.startswith('pdf')
        aspect_ratio = format.split('_')[1]
        
        if is_pdf:
            raise HTTPException(status_code=501, detail="PDF export not yet implemented")
        
        # Extract slides from the project
        project_path = Path(__file__).parent.parent.parent / 'projects' / demo.folder_name
        slides_data = extract_hero_slides_from_project(project_path)
        
        if not slides_data:
            raise HTTPException(status_code=404, detail="No hero slides found in project")
        
        # Create PowerPoint
        output = create_powerpoint_from_slides(slides_data, aspect_ratio)
        filename = f"{demo.folder_name}_export.pptx"
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        return StreamingResponse(
            io.BytesIO(output),
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

