from fastapi import APIRouter, HTTPException, Depends, Query
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from db.session import get_db
from models.demonstration import Demonstration
from core.process_manager import process_manager
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
from pathlib import Path
import io
import json
import re
from typing import Literal, List, Dict, Any

router = APIRouter()

# This is the slide data structure from the hero section
SLIDE_DATA = [
    {
        "title": "TQM Digitalization & Integration",
        "subtitle": "Modern • Interactive • SVG-driven",
        "description": "Cover: Total Quality Management digitalisation and integration, with modern visuals and motion.",
        "colors": "from-blue-600 to-purple-700",
        "points": []
    },
    {
        "title": "Objectives · Digitalization Focus",
        "subtitle": "Unified · Accurate · Automated · Accessible · AI-Ready",
        "description": "Five focus areas for a digital TQM stack.",
        "colors": "from-indigo-600 to-cyan-600",
        "points": [
            "Integration of Information Flow — One unified data backbone (PostgreSQL + APIs + event bus).",
            "Zero Data Error / Accuracy — Validation, traceability, audit trails, versioning, auto-sync.",
            "Automation to Reduce Human Error — Workflows + Celery for approvals, SPC alerts, docs, calibration.",
            "Decision-maker Access — Metabase BI, KPI cockpit, drill-down analytics.",
            "AI-Powered Local Tools — On-prem AI (Python/ONNX/CoreML) for anomaly, AOI, prediction, NLP."
        ]
    },
    {
        "title": "Key System Pillars",
        "subtitle": "3 Foundational Pillars",
        "description": "Integrated Data • Workflow Automation • AI-Assisted Operations",
        "colors": "from-fuchsia-600 to-rose-600",
        "points": []
    },
    {
        "title": "Integrated Data Core",
        "subtitle": "Secure · Consistent · Exemplified",
        "description": "Unified schema, access control, and reliable cross-module synchronization.",
        "colors": "from-emerald-600 to-teal-600",
        "points": [
            "PostgreSQL unified schema (core, project, quality, mfg, supplier, docs, ci, wf).",
            "RLS + ABAC for user-level visibility and control.",
            "Change-data capture → Redis Streams for sync.",
            "Data Validation Layer (JSON Schema + SQL constraints + ETL rules)."
        ]
    },
    {
        "title": "Workflow & Approval Automation",
        "subtitle": "Compliance · Traceability · Speed",
        "description": "Central engine for states, approvals, notifications, and SLAs.",
        "colors": "from-amber-600 to-orange-600",
        "points": [
            "Central workflow engine + approval engine (Part 11-style e-signs).",
            "ていto-notifications, SLA tracking, escalation.",
            "Automations: SPC violation → CAPA; CAPA closed → CI follow-up; New SOP → training tasks."
        ]
    },
    {
        "title": "AI-Assisted Operations",
        "subtitle": "Local · Private · Insightful",
        "description": "On-prem AI models for manufacturing and quality intelligence.",
        "colors": "from-violet-600 to-blue-600",
        "points": [
            "AOI image analysis (defects, features).",
            "SPC drift prediction via local ML.",
            "NLP document assistant for key-field extraction.",
            "Chat interface: ask \"PPM trend for Vendor A this quarter\"."
        ]
    },
    {
        "title": "Decision Support & Analytics",
        "subtitle": "Dashboards · Drill-down · KPIs",
        "description": "Metabase-powered analytics with operational evidence links.",
        "colors": "from-sky-600 to-cyan-600",
        "points": [
            "Panels: Project Health, CAPA aging, NC by cause, Supplier scorecards.",
            "SPC & Yield trends, CI savings, AI anomaly alerts.",
            "Drill-down → root-cause records → document evidence."
        ]
    },
    {
        "title": "Automation & Zero-Error Workflow",
        "subtitle": "Reliable · Deterministic · Auditable",
        "description": "Automated IDs, validations, recalculations, and scheduled tasks.",
        "colors": "from-slate-700 to-gray-800",
        "points": [
            "Auto IDs, timestamps, checksums; file SHA-256 validation.",
            "Auto-KPI recalculation on record updates.",
            "Celery tasks: approvals, calibration alerts, SPC rebuilds, backups & syncs."
        ]
    }
]


def parse_gradient_colors(gradient_str: str) -> tuple:
    """Parse Tailwind gradient classes to get RGB colors"""
    # Default colors
    default_colors = {
        'from-blue-600 to-purple-700': (RGBColor(37, 99, 235), RGBColor(126, 34, 206)),
        'from-indigo-600 to-cyan-600': (RGBColor(79, 70, 229), RGBColor(8, 145, 178)),
        'from-fuchsia-600 to-rose-600': (RGBColor(192, 38, 211), RGBColor(225, 29, 72)),
        'from-emerald-600 to-teal-600': (RGBColor(5, 150, 105), RGBColor(13, 148, 136)),
        'from-amber-600 to-orange-600': (RGBColor(217, 119, 6), RGBColor(234, 88, 12)),
        'from-violet-600 to-blue-600': (RGBColor(124, 58, 237), RGBColor(37, 99, 235)),
        'from-sky-600 to-cyan-600': (RGBColor(2, 132, 199), RGBColor(8, 145, 178)),
        'from-slate-700 to-gray-800': (RGBColor(51, 65, 85), RGBColor(31, 41, 55)),
    }
    
    return default_colors.get(gradient_str, (RGBColor(59, 130, 246), RGBColor(147, 51, 234)))


def create_powerpoint_from_data(slides_data: List[Dict], aspect_ratio: str) -> bytes:
    """Create PowerPoint with actual text boxes and shapes from slide data"""
    prs = Presentation()
    
    if aspect_ratio == '16:9':
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    for data in slides_data:
        # Create blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Set gradient background
        colors = parse_gradient_colors(data.get('colors', ''))
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors[0]  # Use first color as background
        
        # Add title
        if data.get('title'):
            left = Inches(1)
            top = Inches(1.5)
            width = prs.slide_width - Inches(2)
            height = Inches(1.5)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = data['title']
            
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(56)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add subtitle
        if data.get('subtitle'):
            left = Inches(1)
            top = Inches(2.8)
            width = prs.slide_width - Inches(2)
            height = Inches(0.8)
            
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = data['subtitle']
            
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(240, 240, 240)
        
        # Add description
        if data.get('description'):
            left = Inches(1)
            top = Inches(3.6)
            width = prs.slide_width - Inches(2)
            height = Inches(1)
            
            desc_box = slide.shapes.add_textbox(left, top, width, height)
            desc_frame = desc_box.text_frame
            desc_frame.text = data['description']
            
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(16)
            desc_para.font.color.rgb = RGBColor(220, 220, 220)
        
        # Add bullet points
        if data.get('points'):
            left = Inches(1.5)
            top = Inches(4.8)
            width = prs.slide_width - Inches(2.5)
            max_height = prs.slide_height - Inches(4.8)
            
            points_box = slide.shapes.add_textbox(left, top, width, max_height)
            points_frame = points_box.text_frame
            points_frame.word_wrap = True
            
            for i, point in enumerate(data['points']):
                if i > 0:
                    points_frame.add_paragraph()
                para = points_frame.paragraphs[i]
                # Clean up the text
                para.text = point.replace('—', '-').replace('•', '').strip()
                para.level = 0
                para.font.size = Pt(13)
                para.font.color.rgb = RGBColor(240, 240, 240)
                para.space_after = Pt(3)
        
        # Add decorative shape
        left = Inches(prs.slide_width.inches - 3.5)
        top = Inches(1.5)
        width = Inches(2.5)
        height = Inches(2.5)
        
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.fill.transparency = 0.15
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(2)
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


@router.post("/export/{demo_id}")
async def export_demo(
    demo_id: int,
    format: Literal['ppt_169', 'ppt_43', 'pdf_169', 'pdf_43'] = Query(...),
    db: Session = Depends(get_db)
):
    """Export demo to PowerPoint or PDF using direct data structure"""
    demo = db.query(Demonstration).filter(Demonstration.id == demo_id).first()
    if not demo:
        raise HTTPException(status_code=404, detail="Demo not found")
    
    # Check if this is quality-implementation
    if demo.folder_name != 'quality-implementation':
        raise HTTPException(status_code=501, detail=f"Export not yet implemented for {demo.folder_name}")
    
    try:
        is_pdf = format.startswith('pdf')
        aspect_ratio = format.split('_')[1]
        
        if is_pdf:
            raise HTTPException(status_code=501, detail="PDF export not yet implemented")
        
        # Create PowerPoint directly from slide data
        output = create_powerpoint_from_data(SLIDE_DATA, aspect_ratio)
        filename = f"{demo.folder_name}_export.pptx"
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        return StreamingResponse(
            io.BytesIO(output),
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

